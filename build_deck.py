#!/usr/bin/env python3
"""
build_deck.py -- the platform's single entry point.

Given a lesson (either a raw Teacher Guide text file, or an already-
extracted lesson JSON matching extract_lesson.py's schema), runs the full
pipeline and produces a finished .pptx:

  1. Extract structured fields + classify lesson type (skipped if a .json
     is passed directly -- useful for testing without a live API key)
  2. Generate per-word vocabulary icons dynamically via icon_resolver.py
  3. Generate the branded cover slide with this lesson's real Grade/Unit
  4. Dispatch to the correct generate_deck*.js for the detected lesson type
  5. Run the required QA checks (content placeholders, file validation)

Usage:
    python3 build_deck.py lesson.txt output.pptx
    python3 build_deck.py lesson.json output.pptx      # skips extraction

Requires ANTHROPIC_API_KEY in the environment when given a .txt lesson.
Not required when given a pre-extracted .json (e.g. for testing).
"""
import sys
import os
import json
import subprocess
import shutil

HERE = os.path.dirname(os.path.abspath(__file__))

GENERATORS = {
    "Reading": os.path.join(HERE, "generate_deck.js"),
    "Close Reading": os.path.join(HERE, "generate_deck.js"),
    "Literature Response": os.path.join(HERE, "generate_deck_literature_response.js"),
}


def log(msg):
    print(f"[build_deck] {msg}", file=sys.stderr)


def maybe_upload_to_drive(file_path, filename):
    """
    Optional: uploads the finished deck to a Google Drive folder using a
    service account, if GOOGLE_SERVICE_ACCOUNT_JSON and
    GOOGLE_DRIVE_FOLDER_ID are both set in the environment. A service
    account (not interactive OAuth) is the right fit here specifically
    because this runs on a server with no human present to click through
    a consent screen -- the account authenticates itself directly, using
    access granted ahead of time by sharing a real Drive folder with the
    service account's own email address.

    Returns a webViewLink (a real, clickable Drive URL) on success, or
    None if the feature isn't configured. A configured-but-failing
    upload logs the real error and returns None rather than raising --
    the deck itself already generated successfully by the time this
    runs, and a Drive misconfiguration shouldn't take that away.
    """
    sa_json = os.environ.get("GOOGLE_SERVICE_ACCOUNT_JSON")
    folder_id = os.environ.get("GOOGLE_DRIVE_FOLDER_ID")
    if not sa_json or not folder_id:
        return None

    try:
        from google.oauth2 import service_account
        from googleapiclient.discovery import build as build_drive_service
        from googleapiclient.http import MediaFileUpload
    except ImportError:
        log("GOOGLE_SERVICE_ACCOUNT_JSON is set, but the Drive libraries "
            "aren't installed (see requirements.txt) -- skipping upload.")
        return None

    try:
        creds = service_account.Credentials.from_service_account_info(
            json.loads(sa_json),
            scopes=["https://www.googleapis.com/auth/drive.file"],
        )
        service = build_drive_service("drive", "v3", credentials=creds)
        media = MediaFileUpload(
            file_path,
            mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        uploaded = service.files().create(
            body={"name": filename, "parents": [folder_id]},
            media_body=media,
            fields="id,webViewLink",
        ).execute()
        link = uploaded.get("webViewLink")
        log(f"Uploaded to Google Drive: {link}")
        return link
    except Exception as e:
        log(f"Google Drive upload failed (deck still generated successfully): {e}")
        return None


def run_step(cmd, step_name):
    """
    Runs a subprocess step and, on failure, raises with the actual
    stdout/stderr captured directly in the exception message -- not just
    printed to the console, where it's easy to lose in log scrollback or
    a truncated copy-paste. Every subprocess call in this pipeline goes
    through this rather than a bare subprocess.run(check=True).
    """
    result = subprocess.run(cmd, capture_output=True, text=True)
    if result.returncode != 0:
        raise SystemExit(
            f"{step_name} failed (exit code {result.returncode}).\n"
            f"Command: {' '.join(cmd)}\n"
            f"--- stdout ---\n{result.stdout.strip()}\n"
            f"--- stderr ---\n{result.stderr.strip()}"
        )
    if result.stdout.strip():
        log(result.stdout.strip())
    return result


def run_extraction(lesson_path):
    sys.path.insert(0, HERE)
    from extract_lesson import extract_from_path
    log("Calling Claude to extract structured lesson fields...")
    return extract_from_path(lesson_path)


def build(lesson_input_path, out_pptx_path, keep_temp=False):
    # ----- Step 1: get the extracted lesson JSON -----
    if lesson_input_path.endswith(".json"):
        log(f"Input is already-extracted JSON: {lesson_input_path}")
        lesson = json.load(open(lesson_input_path, encoding="utf-8"))
    else:
        lesson = run_extraction(lesson_input_path)

    lesson_type = lesson.get("lesson_type")
    if not lesson.get("supported", True) or lesson_type not in GENERATORS:
        raise SystemExit(
            f"Lesson type {lesson_type!r} is not yet supported by this "
            f"platform. Currently supported: {list(GENERATORS.keys())}. "
            f"No slide deck was generated -- building a generator for this "
            f"type is separate work, not something to fake."
        )
    generator_script = GENERATORS[lesson_type]
    log(f"Detected lesson_type={lesson_type!r} -> using {os.path.basename(generator_script)}")

    work_dir = out_pptx_path + "_build_tmp"
    os.makedirs(work_dir, exist_ok=True)
    try:
        # ----- Step 2: generate vocabulary icons dynamically -----
        lesson_json_path = os.path.join(work_dir, "lesson.json")
        json.dump(lesson, open(lesson_json_path, "w", encoding="utf-8"), indent=2)

        icons_dir = os.path.join(work_dir, "icons")
        log("Generating vocabulary icons...")
        run_step(
            ["python3", os.path.join(HERE, "make_vocab_icons.py"), lesson_json_path, icons_dir],
            "Vocabulary icon generation",
        )

        # ----- Step 3: generate the branded cover slide -----
        grade = lesson.get("grade", "").replace("Grade ", "")
        unit = lesson.get("unit_number", "")
        cover_text = f"Grade {grade}, Knowledge Unit {unit}"
        cover_png = os.path.join(work_dir, "cover.png")
        log(f"Generating cover slide for {cover_text!r}...")
        run_step(
            ["python3", os.path.join(HERE, "assets", "make_cover.py"), cover_text, cover_png],
            "Cover slide generation",
        )

        # ----- Step 4: generate the deck -----
        log("Generating slide deck...")
        run_step(
            ["node", generator_script, lesson_json_path, out_pptx_path, cover_png, icons_dir],
            "Slide deck generation",
        )

        # ----- Step 5: QA -----
        # Uses only python-pptx (a standard pip package, listed in
        # requirements.txt) rather than this sandbox's internal skill
        # scripts, so QA actually runs the same way on a real server as
        # it does here.
        log("Running QA checks...")
        try:
            from pptx import Presentation
            prs = Presentation(out_pptx_path)
            slide_count = len(prs.slides)
            if slide_count == 0:
                raise SystemExit("QA failed: generated file has zero slides.")
            log(f"File opens correctly, {slide_count} slides.")

            leftover = []
            for i, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    text = shape.text_frame.text
                    if "[insert" in text.lower() or "[title]" in text.lower():
                        leftover.append(f"slide {i}: {text[:60]!r}")
            if leftover:
                log(f"WARNING: {len(leftover)} unfilled placeholder(s) remain: {leftover[:5]}")
            else:
                log("No leftover template placeholders found.")
        except SystemExit:
            raise
        except Exception:
            raise SystemExit(f"QA validation failed -- could not open generated file:\n{traceback.format_exc()}")

        # Always preserve a copy of what was actually extracted, right
        # next to the final output -- regardless of keep_temp. This is
        # the only way to answer "why did the deck come out looking like
        # X" questions with real data instead of guessing: the extracted
        # JSON is the ground truth for what Claude actually returned,
        # independent of whether the generator rendered it as intended.
        extracted_json_path = out_pptx_path + ".extracted.json"
        shutil.copy(lesson_json_path, extracted_json_path)
        log(f"Saved extracted lesson JSON to: {extracted_json_path}")

        drive_link = maybe_upload_to_drive(out_pptx_path, os.path.basename(out_pptx_path))

        log(f"Done: {out_pptx_path}")
        return {"drive_link": drive_link}
    finally:
        if not keep_temp:
            shutil.rmtree(work_dir, ignore_errors=True)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        raise SystemExit("Usage: python3 build_deck.py <lesson.txt|lesson.json> <output.pptx>")
    build(sys.argv[1], sys.argv[2], keep_temp="--keep-temp" in sys.argv)
